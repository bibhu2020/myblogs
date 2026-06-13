"""
Meridian — System Design PPTX Generator
Produces meridian_sysdesign.pptx

Slides:
  0  Title
  1  Introduction — Why Meridian Exists & Key Features
  2  Platform Highlights (feature card grid)
  3  System Architecture (diagram + bullets)
  4  Architecture Deep Dive (component table / text)
  5–18  Sequence Diagram slides ×14 (diagram left, explanation right)
  19 Technology Stack
  20 Future Improvements
  21 Thank You
"""
import io, os, textwrap
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ────────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0f, 0x17, 0x2a)
BG_MID    = RGBColor(0x16, 0x22, 0x36)
BG_CARD   = RGBColor(0x1e, 0x29, 0x47)
BG_LIGHT  = RGBColor(0xf4, 0xf6, 0xfa)
BG_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_VIOLET  = RGBColor(0x7c, 0x3a, 0xed)
C_VLIGHT  = RGBColor(0xc4, 0xb5, 0xfd)
C_TEAL    = RGBColor(0x06, 0xb6, 0xd4)
C_GREEN   = RGBColor(0x34, 0xd3, 0x99)
C_AMBER   = RGBColor(0xfb, 0xbf, 0x24)
C_ROSE    = RGBColor(0xfb, 0x71, 0x85)
C_WHITE   = RGBColor(0xff, 0xff, 0xff)
C_GRAY    = RGBColor(0x94, 0xa3, 0xb8)
C_DARK    = RGBColor(0x0d, 0x0d, 0x1a)

# mpl hex strings
VIOLET = '#7c3aed'; BLUE = '#2563eb'; TEAL = '#0891b2'
GREEN  = '#16a34a'; AMBER = '#d97706'; RED  = '#dc2626'
SLATE  = '#334155'; ROSE  = '#e11d48'; DARK = '#0d0d1a'

# ── Presentation setup ─────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# ── Low-level drawing helpers ──────────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill: RGBColor,
             border: RGBColor = None, bw: float = 1.0):
    shp = slide.shapes.add_shape(1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(bw)
    else:
        shp.line.fill.background()
    return shp

def add_text(slide, text, x, y, w, h,
             size=12, color: RGBColor = C_WHITE,
             bold=False, italic=False,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txb

def add_para(tf, text, size=10, color: RGBColor = C_WHITE,
             bold=False, space_before=3, align=PP_ALIGN.LEFT):
    try:
        p = tf.add_paragraph()
    except Exception:
        p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p

def add_image_bytes(slide, png_bytes, x, y, w, h):
    slide.shapes.add_picture(
        io.BytesIO(png_bytes),
        Inches(x), Inches(y), Inches(w), Inches(h))

# ── Composite helpers ──────────────────────────────────────────────────────────
def slide_header(slide, title, subtitle=None, dark=True):
    """Full-width top bar."""
    bg_col = BG_MID if dark else C_VIOLET
    add_rect(slide, 0, 0, 13.33, 0.68, bg_col)
    add_text(slide, title, 0.25, 0.05, 10, 0.55,
             size=20, bold=True, color=C_VLIGHT if dark else C_WHITE)
    if subtitle:
        add_text(slide, subtitle, 10.5, 0.1, 2.6, 0.5,
                 size=9, color=C_GRAY if dark else C_WHITE,
                 italic=True, align=PP_ALIGN.RIGHT)

def bullet_panel(slide, title, bullets,
                 x=8.72, y=0.72, w=4.38, h=6.6,
                 title_color=C_VLIGHT):
    """Dark right-hand explanation card."""
    add_rect(slide, x, y, w, h, BG_CARD, C_VIOLET, 1.2)
    # Title bar inside card
    add_rect(slide, x, y, w, 0.46, BG_MID)
    add_text(slide, title, x+0.12, y+0.05, w-0.18, 0.38,
             size=10.5, bold=True, color=title_color)

    # Bullets
    txb = slide.shapes.add_textbox(
        Inches(x+0.14), Inches(y+0.56),
        Inches(w-0.22), Inches(h-0.68))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4 if i else 0)
        run = p.add_run()
        run.text = f"▸  {bullet}"
        run.font.size = Pt(9)
        run.font.color.rgb = C_WHITE

def feature_card(slide, x, y, w, h, icon, title, desc,
                 accent: RGBColor = C_VIOLET):
    add_rect(slide, x, y, w, h, BG_CARD, accent, 1.0)
    add_text(slide, icon, x+0.15, y+0.12, 0.5, 0.45, size=18)
    add_text(slide, title, x+0.7, y+0.14, w-0.85, 0.35,
             size=11, bold=True, color=C_VLIGHT)
    add_text(slide, desc, x+0.15, y+0.6, w-0.3, h-0.7,
             size=8.5, color=C_GRAY, wrap=True)


# ── Matplotlib helpers (reused from PDF generator) ─────────────────────────────

def render_sequence(title, participants, steps, figsize=(11, None)):
    n = len(participants)
    line_h = 0.55
    top_pad, bot_pad = 1.4, 0.6
    total_h = top_pad + len(steps) * line_h + bot_pad
    if figsize[1] is None:
        figsize = (figsize[0], max(4, total_h))

    fig, ax = plt.subplots(figsize=figsize)
    ax.set_xlim(-0.55, n - 0.45)
    ax.set_ylim(0, total_h)
    ax.axis('off')
    fig.patch.set_facecolor('#f4f6fa')
    ax.set_facecolor('#f4f6fa')

    base_colors = [VIOLET, BLUE, TEAL, GREEN, AMBER, RED, SLATE, ROSE]
    cmap = {i: base_colors[i % len(base_colors)] for i in range(n)}

    def xp(i): return i

    ax.text((n-1)/2, total_h-0.25, title,
            ha='center', va='top', fontsize=11, fontweight='bold',
            color=DARK, fontfamily='DejaVu Sans')

    box_y_top = total_h - 0.55
    box_h = 0.52
    for i, name in enumerate(participants):
        col = cmap[i]
        ax.add_patch(mpatches.FancyBboxPatch(
            (xp(i)-0.35, box_y_top-box_h), 0.70, box_h,
            boxstyle='round,pad=0.05', lw=1.2,
            edgecolor=col, facecolor=col))
        ax.text(xp(i), box_y_top-box_h/2,
                '\n'.join(textwrap.wrap(name, 11)),
                ha='center', va='center', fontsize=7, fontweight='bold',
                color='white', fontfamily='DejaVu Sans')

    lifeline_top = box_y_top - box_h
    lifeline_bot = bot_pad - 0.1
    for i in range(n):
        ax.plot([xp(i), xp(i)], [lifeline_bot, lifeline_top],
                color=cmap[i], lw=0.9, alpha=0.3, ls='--', zorder=1)

    y_cur = lifeline_top - line_h * 0.5
    for step in steps:
        frm, to, lbl, style = step
        y = y_cur
        if style == 'note':
            nx = xp(to)+0.05 if to < n-1 else xp(to)-0.72
            ax.add_patch(mpatches.FancyBboxPatch(
                (nx, y-0.14), 0.72, 0.28,
                boxstyle='round,pad=0.03', lw=0.7,
                edgecolor=AMBER, facecolor='#fef3c7', zorder=4))
            ax.text(nx+0.36, y, lbl,
                    ha='center', va='center', fontsize=6,
                    color='#92400e', zorder=5, style='italic')
        else:
            xi, xo = xp(frm), xp(to)
            col = cmap[frm if style != 'return' else to]
            ls  = '--' if style in ('return','async') else '-'
            ax.annotate('', xy=(xo, y), xytext=(xi, y),
                arrowprops=dict(arrowstyle='->', color=col,
                    lw=1.3 if style=='sync' else 0.9,
                    linestyle=ls,
                    connectionstyle='arc3,rad=0.0'), zorder=3)
            mid = (xi+xo)/2
            ax.text(mid, y+0.08,
                    lbl if len(lbl)<=38 else lbl[:35]+'…',
                    ha='center', va='bottom', fontsize=6.5,
                    color=DARK,
                    bbox=dict(boxstyle='round,pad=0.12', fc='white',
                              ec='none', alpha=0.88),
                    zorder=5)
        y_cur -= line_h

    # Bottom mirrors
    bot_y = lifeline_bot - box_h + 0.1
    for i, name in enumerate(participants):
        col = cmap[i]
        ax.add_patch(mpatches.FancyBboxPatch(
            (xp(i)-0.35, bot_y), 0.70, box_h,
            boxstyle='round,pad=0.05', lw=1.2,
            edgecolor=col, facecolor=col))
        ax.text(xp(i), bot_y+box_h/2,
                '\n'.join(textwrap.wrap(name, 11)),
                ha='center', va='center', fontsize=7, fontweight='bold',
                color='white')

    plt.tight_layout(pad=0.2)
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='#f4f6fa')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def render_arch():
    fig, ax = plt.subplots(figsize=(13, 7.5))
    ax.set_xlim(0, 13); ax.set_ylim(0, 7.5)
    ax.axis('off')
    fig.patch.set_facecolor('#f4f6fa')

    def box(x, y, w, h, label, sub='', col=VIOLET, tc='white', fs=8.5):
        ax.add_patch(mpatches.FancyBboxPatch(
            (x,y), w, h, boxstyle='round,pad=0.1',
            lw=1.4, edgecolor=col, facecolor=col))
        ax.text(x+w/2, y+h/2+(0.16 if sub else 0), label,
                ha='center', va='center', fontsize=fs,
                fontweight='bold', color=tc)
        if sub:
            ax.text(x+w/2, y+h/2-0.2, sub,
                    ha='center', va='center', fontsize=6.5,
                    color=tc, alpha=0.85, style='italic')

    def arr(x1, y1, x2, y2, lbl='', col=SLATE):
        ax.annotate('', xy=(x2,y2), xytext=(x1,y1),
            arrowprops=dict(arrowstyle='->', color=col, lw=1.4,
                            connectionstyle='arc3,rad=0.0'))
        if lbl:
            ax.text((x1+x2)/2, (y1+y2)/2+0.1, lbl,
                    ha='center', va='bottom', fontsize=6.5, color=SLATE)

    def dash(x1, y1, x2, y2, lbl='', col='#64748b'):
        ax.annotate('', xy=(x2,y2), xytext=(x1,y1),
            arrowprops=dict(arrowstyle='->', color=col, lw=1.1,
                            linestyle='dashed',
                            connectionstyle='arc3,rad=0.0'))
        if lbl:
            ax.text((x1+x2)/2, (y1+y2)/2+0.1, lbl,
                    ha='center', va='bottom', fontsize=6.5, color=col, style='italic')

    ax.text(6.5, 7.3, 'Meridian — System Architecture',
            ha='center', va='top', fontsize=13, fontweight='bold', color=DARK)

    # Browser
    box(0.2, 5.4, 2.3, 1.0, 'Browser', 'HF Space / localhost', BLUE)
    # Vite / Nginx
    box(0.2, 3.7, 2.3, 0.95, 'Vite / Nginx', 'Vue 3 SPA', '#0284c7', fs=8)
    arr(1.35, 5.4, 1.35, 4.65)

    # API Gateway
    box(3.5, 4.25, 2.4, 1.1, 'API Gateway', ':3000 · JWT Guard', VIOLET)
    arr(2.5, 5.0, 3.5, 4.9, '/api/*')
    arr(2.5, 4.1, 3.5, 4.5, '/uploads/*')

    # Auth Service
    box(6.8, 5.8, 2.3, 0.95, 'Auth Service', ':3001 · auth.db', TEAL)
    arr(5.9, 5.0, 6.8, 6.2, '/api/auth/*')

    # Blog Service
    box(6.8, 4.15, 2.3, 0.95, 'Blog Service', ':3002 · blog.db', GREEN)
    arr(5.9, 4.75, 6.8, 4.6, '/api/posts/*')

    # Media Service
    box(6.8, 2.5, 2.3, 0.95, 'Media Service', ':3003 · media.db', AMBER)
    arr(5.9, 4.4, 6.8, 2.95, '/api/media/*')

    # TTS Service
    box(6.8, 0.85, 2.3, 0.95, 'TTS Service', ':5500 · Kokoro AI', ROSE)
    arr(5.9, 4.25, 6.8, 1.3, '/api/tts')

    # GitHub
    box(9.8, 2.5, 2.5, 0.95, 'GitHub', 'bibhu2020/media CDN', '#1f2937', fs=8)
    dash(9.1, 2.95, 9.8, 2.95, 'PUT image')

    # OpenAI
    box(9.8, 5.1, 2.5, 0.95, 'OpenAI', 'GPT-4o / DALL-E', '#065f46', fs=8)
    # HF Inference
    box(9.8, 3.75, 2.5, 0.95, 'HF Inference', 'FLUX.1-schnell', '#b45309', tc='white', fs=8)
    # Gemini
    box(9.8, 6.3, 2.5, 0.85, 'Gemini', '2.5-flash', '#1d4ed8', fs=8)
    # Unsplash
    box(9.8, 1.35, 2.5, 0.85, 'Unsplash', 'Photo fallback', '#374151', fs=8)

    # AI Agent
    box(3.5, 6.2, 2.4, 1.05, 'AI Agent', 'python-pptx / node-cron', '#7e22ce')
    arr(4.7, 6.2, 4.7, 5.35, 'POST /api/mcp')
    dash(5.9, 6.6, 9.8, 6.55, 'research + image gen')

    # Supervisord label
    ax.text(6.5, 0.22,
            'All services run in a single Docker container managed by Supervisord',
            ha='center', va='bottom', fontsize=7, color='#6b7280', style='italic')

    # Legend
    items = [(VIOLET,'API GW'),(TEAL,'Auth'),(GREEN,'Blog'),
             (AMBER,'Media'),(ROSE,'TTS'),(BLUE,'Frontend'),('#7e22ce','AI Agent')]
    for idx,(col,lbl) in enumerate(items):
        px = 0.3 + idx*1.82
        ax.add_patch(mpatches.FancyBboxPatch(
            (px,0.05),0.22,0.2,boxstyle='round,pad=0.03',
            fc=col, ec='none'))
        ax.text(px+0.28, 0.15, lbl, va='center', fontsize=6.5, color=DARK)

    plt.tight_layout(pad=0.2)
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='#f4f6fa')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ── Flow definitions: (title, participants, steps, bullets) ────────────────────

FLOWS = [

("User Registration",
 ['Browser','Vue SPA','API Gateway\n:3000','Auth Service\n:3001','auth.db'],
 [(0,1,'Navigate to /admin','sync'),
  (1,1,'router.beforeEach → not logged in','note'),
  (1,0,'Redirect → /admin/login','return'),
  (0,1,'POST /api/auth/register { name, email, password }','sync'),
  (1,2,'POST /api/auth/register','sync'),
  (2,3,'POST /auth/register (no JWT required)','sync'),
  (3,3,'bcrypt.hash(password, 10)','note'),
  (3,4,'INSERT INTO users','sync'),
  (4,3,'201 { id, email }','return'),
  (3,2,'201 Created','return'),
  (2,1,'201 Created','return'),
  (1,0,'Prompt login — account created','return')],
 ["Browser navigates to /admin; Vue Router guard detects no token",
  "Guard redirects to /admin/login (no API call needed)",
  "Registration form POSTs through API Gateway — no JWT required",
  "Gateway forwards to Auth Service :3001",
  "Auth Service bcrypt-hashes the password (10 salt rounds)",
  "Hashed credentials written to auth.db users table",
  "201 Created bubbles back; user prompted to log in"]),

("User Login & JWT Issue",
 ['Browser','Vue SPA\nauth.js store','API Gateway\n:3000','Auth Service\n:3001'],
 [(0,1,'POST /api/auth/login { email, password }','sync'),
  (1,2,'POST /api/auth/login','sync'),
  (2,3,'POST /auth/login','sync'),
  (3,3,'bcrypt.compare(pw, hash)','note'),
  (3,2,'200 { access_token, user }','return'),
  (2,1,'200 { access_token, user }','return'),
  (1,1,'localStorage.setItem("token", jwt)','note'),
  (1,1,'axios.defaults.headers.Authorization = Bearer…','note'),
  (1,0,'Redirect → /admin/dashboard','return'),
  (0,1,'All subsequent calls include Authorization header','async'),
  (1,2,'Gateway validates JWT via JwtStrategy','async')],
 ["Credentials POST through Gateway; Auth Service validates with bcrypt",
  "On match, Auth Service signs a JWT (HS256, 24h TTL)",
  "Token stored in localStorage by Pinia auth store",
  "Axios default header set — all subsequent requests carry the token",
  "Browser redirected to /admin/dashboard",
  "Every protected route is gated at the Gateway — internal services re-verify the forwarded header"]),

("A/B Layout Assignment",
 ['Browser','Vue SPA\nlayout store','localStorage','Pinia Store','Page\nComponents'],
 [(0,1,'First visit — no variant in localStorage','sync'),
  (1,2,'getItem("meridian_ab_variant")','sync'),
  (2,1,'null (not set)','return'),
  (1,1,'Math.random() < 0.5 → "a" or "b"','note'),
  (1,2,'setItem("meridian_ab_variant", "a"|"b")','sync'),
  (1,3,'variant = ref("a"|"b")','sync'),
  (3,4,'useLayoutStore() injects variant','sync'),
  (4,4,'layout.variant==="b" → dark classes','note'),
  (0,1,'Return visit — variant already stored','async'),
  (1,2,'getItem(…) → "a" or "b" (sticky)','async'),
  (1,3,'Same variant re-used — no re-randomisation','sync')],
 ["On first visit, localStorage has no variant — coin-flip assignment",
  "Variant A = Editorial Light; B = Dark Magazine",
  "Result persisted to localStorage so same user always sees same layout",
  "Pinia layout store exposes variant as a reactive ref",
  "Every component uses :class or v-if bindings against layout.variant",
  "Return visits read the stored value — session-sticky by design"]),

("Browse Home / All Posts",
 ['Browser','Vue SPA\nblog store','API Gateway\n:3000','Blog Service\n:3002','blog.db'],
 [(0,1,'Navigate to / or /blog','sync'),
  (1,1,'blog.fetchCategories() + fetchPosts()','note'),
  (1,2,'GET /api/categories','sync'),
  (2,3,'GET /categories','sync'),
  (3,4,'SELECT * FROM categories','sync'),
  (4,3,'rows[]','return'),
  (3,2,'200 [ { id, name, slug, icon } ]','return'),
  (2,1,'categories[]','return'),
  (1,2,'GET /api/posts?page=1&limit=12','sync'),
  (2,3,'GET /posts?page=1&limit=12','sync'),
  (3,4,'SELECT … ORDER BY createdAt DESC','sync'),
  (4,3,'posts[], total','return'),
  (3,2,'200 { posts[], pagination }','return'),
  (2,1,'posts[], pagination','return'),
  (1,0,'Render PostCard grid (layout A or B)','return')],
 ["Home & Blog pages trigger two lazy fetches via Pinia blog store",
  "Categories and posts fetched concurrently",
  "Posts paginated (12 per page), sorted by creation date DESC",
  "Pinia caches results — navigating back skips re-fetch",
  "PostCard grid renders in A (light editorial) or B (dark magazine) layout",
  "Category nav, tag chips, and pagination all driven from cached store state"]),

("Read Individual Blog Post",
 ['Browser','BlogPost.vue','API Gateway\n:3000','Blog Service\n:3002'],
 [(0,1,'Navigate to /blog/:slug','sync'),
  (1,2,'GET /api/posts/:slug','sync'),
  (2,3,'GET /posts/:slug','sync'),
  (3,3,'SELECT post; UPDATE views += 1','note'),
  (3,2,'200 { post, category, tags }','return'),
  (2,1,'post object','return'),
  (1,1,'v-html content; hljs.highlightElement()','note'),
  (1,2,'GET /api/posts?category=slug&limit=3','sync'),
  (2,3,'GET /posts?category=slug&limit=3','sync'),
  (3,2,'200 related posts[]','return'),
  (2,1,'related posts','return'),
  (1,2,'GET /api/comments/post/:id','sync'),
  (2,3,'GET /comments/post/:id','sync'),
  (3,2,'200 comments[]','return'),
  (2,1,'comments[]','return'),
  (1,0,'Full page: article + related + comments','return')],
 ["Slug-based lookup — no numeric IDs exposed in URLs",
  "Blog Service increments the view counter on every fetch",
  "Content rendered via v-html; highlight.js syntax-highlights code blocks post-render",
  "Related posts fetched by same category (limit 3) in a follow-up request",
  "Comments fetched separately; only approved comments shown to readers",
  "All three requests resolve before full page paint"]),

("Text-to-Speech Playback",
 ['User','BlogPost.vue','API Gateway\n:3000','TTS Service\n:5500','Kokoro AI\nModel'],
 [(0,1,'Click "Listen to article"','sync'),
  (1,1,'splitChunks(content) → N chunks ≤500 chars','note'),
  (1,1,'Dispatch all N fetches concurrently','note'),
  (1,2,'POST /api/tts { text: chunk[0] }','sync'),
  (1,2,'POST /api/tts { text: chunk[1] } (parallel)','async'),
  (2,3,'POST /tts { text }','sync'),
  (3,4,'kokoro.generate(text)','sync'),
  (4,3,'WAV audio bytes','return'),
  (3,2,'200 audio/wav blob','return'),
  (2,1,'Blob (chunk 0 ready)','return'),
  (1,1,'new Audio(blobUrl).play()','note'),
  (0,1,'Click Pause / Seek / Stop','sync'),
  (1,1,'sessionId++ → all loops invalidated','note')],
 ["Post content split into ≤500-char semantic chunks",
  "All N chunks dispatched as concurrent fetch requests — pre-fetch pipeline",
  "Playback starts with chunk 0 while later chunks download in background",
  "Kokoro AI is a local on-device model — no external TTS API cost",
  "Seek slider maps fraction → chunk index; cached blobs play immediately",
  "Stop increments a sessionId — all async loops check this and abort cleanly"]),

("Search Articles",
 ['User','Search.vue','API Gateway\n:3000','Blog Service\n:3002'],
 [(0,1,'Type query in search box','sync'),
  (1,1,'350ms debounce timer','note'),
  (1,2,'GET /api/posts?search=query&limit=50','sync'),
  (2,3,'GET /posts?search=query&limit=50','sync'),
  (3,3,'WHERE title LIKE %q% OR content LIKE %q%','note'),
  (3,2,'200 { posts[], total }','return'),
  (2,1,'results[]','return'),
  (1,0,'Render result cards with thumbnail + meta','return'),
  (0,1,'Click category chip in sidebar','sync'),
  (1,2,'GET /api/posts?category=slug&search=query','sync'),
  (2,3,'Filtered & searched results','sync'),
  (3,2,'200 filtered[]','return'),
  (2,1,'results[]','return'),
  (1,0,'Update results + URL (router.replace)','return')],
 ["350ms debounce prevents server flood while typing",
  "Blog Service searches title AND content with SQL LIKE operators",
  "Returns up to 50 results, client renders cards with thumbnail and meta",
  "Category filter chips refine without clearing the search query",
  "URL updated via router.replace() — search links are bookmarkable",
  "Sidebar shows matching category counts for easy drill-down"]),

("Submit Comment",
 ['Reader','BlogPost.vue','API Gateway\n:3000','Blog Service\n:3002','blog.db'],
 [(0,1,'Fill name / email / content, click Post','sync'),
  (1,2,'POST /api/comments/post/:id { authorName, email, content }','sync'),
  (2,3,'POST /comments/post/:id','sync'),
  (3,4,'INSERT INTO comments (status=pending)','sync'),
  (4,3,'comment row','return'),
  (3,2,'201 { id, status:"pending" }','return'),
  (2,1,'201 Created','return'),
  (1,0,'"Comment awaiting approval" banner shown','return'),
  (None,2,'Admin reviews at /admin (PATCH approve)','note'),
  (None,3,'Approved comments become visible to readers','note')],
 ["No authentication required — open commenting for all readers",
  "Comments stored immediately with status='pending'",
  "Pending comments are invisible to public readers",
  "Reader sees a 'awaiting approval' confirmation banner instantly",
  "Admin reviews in /admin panel and PATCHes status to 'approved'",
  "Protects from spam without blocking legitimate engagement"]),

("Admin Create Blog Post",
 ['Admin','PostEditor.vue','API Gateway\n:3000','Blog Service\n:3002','blog.db'],
 [(0,1,'Navigate to /admin/posts/new','sync'),
  (1,1,'JWT guard validates token in localStorage','note'),
  (0,1,'Write content in TipTap rich-text editor','sync'),
  (1,1,'TipTap → HTML with syntax-highlighted code','note'),
  (0,1,'Click Publish','sync'),
  (1,2,'POST /api/posts { title, content, categoryId, tags[], featuredImage }','sync'),
  (2,2,'JwtAuthGuard validates Bearer token','note'),
  (2,3,'POST /posts { …body, authorId, authorName }','sync'),
  (3,4,'INSERT INTO posts; link tags in junction table','sync'),
  (4,3,'post row','return'),
  (3,2,'201 { post }','return'),
  (2,1,'201 Created','return'),
  (1,0,'Redirect → /admin/posts (success toast)','return')],
 ["Route /admin/posts/new protected by Vue Router JWT guard",
  "TipTap editor supports rich HTML: headings, blockquote, code, images, links",
  "CodeBlockLowlight provides syntax highlighting for 30+ languages",
  "JWT validated at Gateway before forwarding — unauthorized requests rejected at the edge",
  "Tags use a junction table posts_tags_tags (many-to-many relationship)",
  "On success, browser redirects to post list with a toast notification"]),

("Media File Upload",
 ['Admin','Vue SPA','API Gateway\n:3000','Media Service\n:3003','GitHub\nbibhu2020/media','media.db'],
 [(0,1,'Select image in post editor','sync'),
  (1,2,'POST /api/media/upload multipart/form-data','sync'),
  (2,2,'JwtAuthGuard validates token','note'),
  (2,3,'forwardWithFile() — streams multipart body','sync'),
  (3,3,'Multer saves with UUID filename to uploads/','note'),
  (3,4,'PUT /repos/bibhu2020/media/contents/:file (base64)','sync'),
  (4,3,'200 OK — raw CDN URL','return'),
  (3,5,'INSERT INTO media { url: github_raw_url }','sync'),
  (5,3,'media row','return'),
  (3,2,'201 { url: "https://raw.githubusercontent.com/…" }','return'),
  (2,1,'media URL','return'),
  (1,1,'Embed URL as featuredImage or in TipTap body','note'),
  (None,3,'GitHub 403 → fallback to /uploads/ local URL','note')],
 ["File uploaded as multipart/form-data through Gateway forwardWithFile()",
  "Multer saves a UUID-named copy to media-service/uploads/ as a safety copy",
  "Media Service base64-encodes and PUTs to GitHub CDN repo via REST API",
  "GitHub returns a permanent raw.githubusercontent.com CDN URL",
  "URL stored in media.db and returned to the editor for embedding",
  "Graceful fallback: if GitHub returns 403, local /uploads/ path is used instead"]),

("AI Agent — Saturday Run (AI Trending)",
 ['node-cron\nSat 3am','Agent\nindex.js','OpenAI\nGPT-4o','Image\nFallback Chain','MCP API\n/api/mcp','Blog\nService'],
 [(0,1,'CRON_SCHEDULE fires (0 3 * * 6)','sync'),
  (1,1,'TOPIC_MODE=ai_trending → AI category','note'),
  (1,2,'discoverTrend() — GPT-4o web_search_preview','sync'),
  (2,1,'trend summary + source URLs','return'),
  (1,2,'deepResearch() — 3 parallel GPT-4o queries','sync'),
  (2,1,'{ technical, reactions, implications }','return'),
  (1,2,'generatePost() — GPT-4o writes full HTML article','sync'),
  (2,1,'{ title, content, excerpt, tags[] }','return'),
  (1,3,'generateImage(prompt) — 5-provider chain','sync'),
  (3,1,'image buffer + credit','return'),
  (1,4,'uploadMedia() then publishPost() via MCP','sync'),
  (4,5,'POST /posts (agent JWT)','sync'),
  (5,4,'201 post created','return'),
  (4,1,'✅ Published','return')],
 ["node-cron fires every Saturday at 3am UTC",
  "TOPIC_MODE=ai_trending → always picks the AI/Technology category",
  "GPT-4o with web_search_preview discovers the week's top AI story",
  "3 parallel research queries add depth: technical analysis, community reactions, future implications",
  "GPT-4o writes a full HTML article (~1500+ words) from combined research context",
  "Image from 5-provider fallback chain (DALL-E 3/2 → FLUX → Gemini → Unsplash)",
  "Post published via MCP REST API using a dedicated agent JWT"]),

("AI Agent — Sunday Run (Random Topic)",
 ['node-cron\nSun 3am','Agent','OpenAI','Image Chain','MCP / Blog'],
 [(0,1,'CRON_SCHEDULE fires (0 3 * * 0)','sync'),
  (1,1,'TOPIC_MODE=random_general','note'),
  (1,1,'pickCategory() — random from 5 categories','note'),
  (1,2,'discoverTrend() with category context','sync'),
  (2,1,'trend + category summary','return'),
  (1,2,'deepResearch() — 3 parallel queries','sync'),
  (2,1,'research object','return'),
  (1,2,'generatePost() — full HTML article','sync'),
  (2,1,'post content','return'),
  (1,3,'generateImage(prompt)','sync'),
  (3,1,'image buffer','return'),
  (1,4,'uploadMedia() + publishPost() via MCP','sync'),
  (4,1,'✅ Published in chosen category','return')],
 ["node-cron fires every Sunday at 3am UTC",
  "TOPIC_MODE=random_general — diversity run",
  "Category randomly chosen from: Technology, Science, History, Travel, Knowledge",
  "Same deep 3-parallel-query research pipeline as Saturday",
  "GPT-4o writes a full article; image generated via fallback chain",
  "Ensures content variety — AI is not the only topic covered",
  "Two runs/week means ~8-9 new posts per month, fully automated"]),

("Image Generation Fallback Chain",
 ['Agent\nimages.js','DALL-E 3\nOpenAI','DALL-E 2\nOpenAI','FLUX.1\nHF Inference','Gemini\nGoogle AI','Unsplash\nAPI'],
 [(0,1,'tryDalle3(prompt, "1792x1024")','sync'),
  (1,0,'Image buffer ✅','return'),
  (None,1,'If error or no access:','note'),
  (0,2,'tryDalle2(prompt, "1024x1024")','sync'),
  (2,0,'Image buffer ✅','return'),
  (None,2,'If error:','note'),
  (0,3,'tryFlux(aiPrompt) — HF Inference API','sync'),
  (3,3,'503: wait estimated_time, retry once','note'),
  (3,0,'Image buffer ✅','return'),
  (None,3,'If timeout / error:','note'),
  (0,4,'tryGemini() — gemini-2.5-flash-image','sync'),
  (4,4,'responseModalities: ["image","text"]','note'),
  (4,0,'Image buffer ✅','return'),
  (None,4,'If no image in response:','note'),
  (0,5,'tryUnsplash(originalPrompt)','sync'),
  (5,0,'Photo URL + credit attribution ✅','return'),
  (None,5,'All fail → null (post without image)','note')],
 ["Five-provider fallback ensures an image is almost always found",
  "DALL-E 3 attempted first — highest resolution (1792×1024) and quality",
  "DALL-E 2 fallback if account lacks DALL-E 3 access or quota exceeded",
  "FLUX.1-schnell via HF Inference API — free tier; smart 503 backoff with one retry",
  "Gemini 2.5 Flash image generation as fourth option",
  "Unsplash keyword search as final safety net — returns photo with attribution credit",
  "If all five fail, post is published without a featured image (very rare)"]),

("CI/CD — GitHub to HF Space Deploy",
 ['Developer','GitHub\nmain','GitHub\nActions','HF Space\ngit repo','Docker\nBuild','Running\nContainer'],
 [(0,1,'git push origin main','sync'),
  (1,2,'Trigger: push to main','sync'),
  (2,2,'Sync 6 secrets to HF Space environment','note'),
  (2,3,'git push hf main (force)','sync'),
  (3,4,'Dockerfile detected → rebuild triggered','sync'),
  (4,4,'npm ci + nest build (4 NestJS services)','note'),
  (4,4,'npm ci + vite build (Vue 3 frontend)','note'),
  (4,4,'pip install (Kokoro TTS deps)','note'),
  (4,5,'supervisord starts all 8 programs','sync'),
  (5,5,'auth, blog, media, gateway, nginx, tts, agents','note'),
  (3,2,'Build complete ✅','return'),
  (2,1,'Workflow completed','return'),
  (1,0,'Live at mishrabp-myblogs.hf.space','return')],
 ["git push to GitHub main triggers GitHub Actions workflow immediately",
  "Action syncs 6 secrets to HF Space: OPENAI, GEMINI, HF_TOKEN, GITHUB, UNSPLASH, SERVER_BASE",
  "Force-pushes code to HuggingFace Space git repository",
  "HF Space detects Dockerfile change and triggers autonomous Docker rebuild",
  "All 4 NestJS services + Vue 3 frontend compiled; Python TTS deps installed",
  "Supervisord starts 8 processes: gateway, auth, blog, media, nginx, tts, 2 agents",
  "Total time: ~30s for git push + ~5min for HF Docker rebuild"]),

]  # end FLOWS


# ── Build presentation ─────────────────────────────────────────────────────────

prs = new_prs()
SW, SH = 13.33, 7.5   # slide dimensions in inches

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 0 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s, BG_DARK)

# Decorative accent bar left
add_rect(s, 0, 0, 0.22, SH, C_VIOLET)
add_rect(s, 0.22, 0, 0.06, SH, BG_MID)

# Main title text
add_text(s, "Meridian", 1.0, 1.3, 9, 1.5,
         size=56, bold=True, color=C_VLIGHT)
add_text(s, "System Design Document", 1.0, 2.85, 9, 0.7,
         size=24, color=C_WHITE)

# Divider
add_rect(s, 1.0, 3.65, 6.5, 0.05, C_VIOLET)

# Subtitle
add_text(s, "A fully AI-built, AI-operated blogging platform", 1.0, 3.8, 10, 0.5,
         size=14, italic=True, color=C_GRAY)

# Metadata grid
meta = [
    ("Stack",      "NestJS ×4 · Vue 3 · SQLite · Tailwind CSS v4"),
    ("AI Features","GPT-4o research · DALL-E / FLUX / Gemini · Kokoro TTS"),
    ("Deployment", "HuggingFace Spaces · Docker · Supervisord · GitHub Actions"),
    ("A/B Testing","Two brand layouts — Editorial Light (A) & Dark Magazine (B)"),
    ("Version",    "June 2026"),
]
for i,(k,v) in enumerate(meta):
    y = 4.4 + i*0.52
    add_rect(s, 1.0, y, 2.1, 0.44, BG_MID)
    add_text(s, k, 1.05, y+0.04, 2.0, 0.36,
             size=9, bold=True, color=C_VLIGHT)
    add_text(s, v, 3.2, y+0.04, 8.5, 0.36,
             size=9, color=C_GRAY)

add_text(s, "Bibhu Mishra · Principal Cloud / DevOps Architect · LTIMindtree",
         1.0, 7.0, 11, 0.38,
         size=9, italic=True, color=C_GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Introduction: What & Why
# ─────────────────────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s, BG_DARK)
add_rect(s, 0, 0, SW, 0.68, BG_MID)
add_text(s, "What is Meridian?", 0.35, 0.06, 8, 0.54,
         size=22, bold=True, color=C_VLIGHT)
add_text(s, "meridian", 10.2, 0.08, 2.9, 0.52,
         size=12, italic=True, color=C_GRAY, align=PP_ALIGN.RIGHT)

# Left column — Why + What
add_rect(s, 0.28, 0.82, 5.1, 1.3, BG_MID, C_VIOLET, 1)
add_text(s, "Why It Exists", 0.45, 0.88, 4.7, 0.38,
         size=11, bold=True, color=C_AMBER)
add_text(s,
    "Bibhu Mishra set out to prove that AI should not just assist developers — it should be the developer. "
    "Meridian is that proof: a production-grade blogging platform where every line of application code was "
    "written by AI, every post is researched and published by AI, and the platform maintains itself with AI.",
    0.42, 1.3, 4.85, 0.78,
    size=8.5, color=C_GRAY, wrap=True)

add_rect(s, 0.28, 2.28, 5.1, 1.18, BG_MID, C_TEAL, 1)
add_text(s, "What Makes It Remarkable", 0.45, 2.34, 4.7, 0.38,
         size=11, bold=True, color=C_TEAL)
points = [
    "Zero human-written application code — 100% Claude AI",
    "2 research-backed posts auto-published every week",
    "Local on-device TTS reads any article aloud (Kokoro AI)",
    "Monthly AI rebranding reflects world events in the UI",
    "AI posts need admin approval before going live (HITL)",
    "MCP API lets any AI agent publish to the platform",
]
txb = s.shapes.add_textbox(Inches(0.44), Inches(2.72), Inches(4.82), Inches(0.72))
tf = txb.text_frame; tf.word_wrap = True
for i, pt in enumerate(points):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.space_before = Pt(1)
    run = p.add_run()
    run.text = f"✦  {pt}"
    run.font.size = Pt(8.2)
    run.font.color.rgb = C_WHITE

# Right column — tag cloud / stats
stats = [
    ("4","NestJS Microservices"),("1","Vue 3 SPA"),("3","SQLite Databases"),
    ("2","AI Agent runs / week"),("5","Image AI providers"),("14","Documented flows"),
    ("2","Layout variants (A/B)"),("∞","Auto-maintenance cycles"),
]
cols = 2
for idx,(num,label_) in enumerate(stats):
    col = idx % cols
    row = idx // cols
    cx = 5.72 + col * 3.76
    cy = 0.82 + row * 1.56
    add_rect(s, cx, cy, 3.52, 1.4, BG_CARD, C_VIOLET, 1)
    add_text(s, num,   cx+0.2,  cy+0.12, 3.0, 0.7,
             size=34, bold=True, color=C_VLIGHT)
    add_text(s, label_,cx+0.2, cy+0.85, 3.0, 0.45,
             size=9,  color=C_GRAY, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Platform Feature Grid
# ─────────────────────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s, BG_DARK)
slide_header(s, "Platform Features at a Glance", "meridian · June 2026")

features = [
    ("🤖","AI-Written Codebase",
     "Every line of application code written by Claude AI through a guided development session. "
     "No human modified application source files."),
    ("📝","Autonomous Publishing",
     "GPT-4o agent runs every Saturday & Sunday at 3am UTC: discovers trending topics, "
     "deep-researches them, writes 1500+ word articles, and publishes with images."),
    ("🔊","Local Text-to-Speech",
     "Kokoro AI model runs on-device (no external API). Any article can be read aloud with "
     "chunked pre-fetch, seek slider, and pause/resume controls."),
    ("🎨","Monthly AI Rebranding",
     "5-agent pipeline (Ideation → Code → Review → Test → Publish) reskins the UI monthly "
     "to reflect world events. Enforces WCAG AA contrast automatically."),
    ("⏸️","Human-in-the-Loop Approval",
     "AI-generated posts are saved as PENDING first. Admin reviews in a dedicated "
     "Pending Approval tab before the post goes live."),
    ("🔌","MCP API for External Agents",
     "Full Model Context Protocol endpoint lets any MCP-compatible AI agent publish, "
     "manage, and query blog content programmatically."),
    ("🪞","A/B Layout Testing",
     "Two complete brand identities — Editorial Light (A) and Dark Magazine (B). "
     "Randomly assigned once per browser; session-sticky via localStorage."),
    ("🛡️","AI-Driven Maintenance",
     "Monthly maintenance agent audits npm CVEs, ADA/WCAG compliance, SEO meta tags, "
     "and open Dependabot PRs — then patches and rebuilds automatically."),
]

cols = 4
for idx, (icon, title_, desc) in enumerate(features):
    col = idx % cols
    row = idx // cols
    cx = 0.22 + col * 3.27
    cy = 0.82 + row * 3.1
    accent = [C_VIOLET, C_TEAL, C_AMBER, C_ROSE,
              C_GREEN,  C_VIOLET, C_TEAL, C_AMBER][idx]
    feature_card(s, cx, cy, 3.1, 2.95, icon, title_, desc, accent)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — System Architecture Diagram
# ─────────────────────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s, BG_LIGHT)
slide_header(s, "System Architecture", "Service topology & external integrations", dark=False)

arch_png = render_arch()
add_image_bytes(s, arch_png, 0.18, 0.74, 8.35, 6.62)

bullet_panel(s,
    "Architecture Overview",
    ["Browser talks to Vite dev server (port 5173) or Nginx in production",
     "All /api/* traffic proxied to API Gateway :3000 — single public entry point",
     "Gateway holds JWT guards; internal services trust the forwarded Authorization header",
     "Auth :3001, Blog :3002, Media :3003 each own an isolated SQLite database",
     "TTS Service runs Kokoro AI as a Python process on :5500",
     "AI Agent publishes via POST /api/mcp using a dedicated service JWT",
     "5-provider image chain: DALL-E 3/2 → FLUX.1 → Gemini → Unsplash",
     "All services run inside one Docker container managed by Supervisord",
     "GitHub Actions CI/CD: push → sync secrets → force-push to HF Space → rebuild"],
    x=8.72, y=0.72, w=4.38, h=6.62)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Architecture Deep Dive
# ─────────────────────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s, BG_DARK)
slide_header(s, "Architecture Deep Dive — Components & Data Model")

# Components table
components = [
    ("API Gateway",    "3000", "—",         "JWT validation; proxies all routes via ProxyService.forward() and forwardWithFile()"),
    ("Auth Service",   "3001", "auth.db",   "Users, bcrypt passwords, JWT sign/validate. Shared secret: myblogs-secret-key-2024"),
    ("Blog Service",   "3002", "blog.db",   "Posts, categories, tags, comments, views. Seeds on empty DB. LIKE full-text search"),
    ("Media Service",  "3003", "media.db",  "Multer upload → GitHub CDN → fallback to local /uploads/. UUID filenames"),
    ("TTS Service",    "5500", "—",         "Python subprocess. POST /tts { text } → WAV blob. Kokoro AI local model"),
    ("AI Agent",       "—",    "—",         "Python (post_agent) + Node (rebranding). Cron-triggered. Publishes via MCP API"),
    ("Vue 3 Frontend", "5173", "—",         "Composition API, Pinia stores, TipTap editor, TTS player, A/B layouts, Tailwind v4"),
    ("Nginx",          "80",   "—",         "Reverse proxy. Routes /api/* → Gateway, /uploads/* → Media Service"),
]

headers = ["Component", "Port", "DB", "Responsibility"]
col_widths = [2.2, 0.7, 1.1, 8.5]  # inches
header_y = 0.82
row_h = 0.54

# Header row
xs = [0.22]
for cw in col_widths[:-1]:
    xs.append(xs[-1] + cw + 0.02)

for j, (hdr, cw) in enumerate(zip(headers, col_widths)):
    add_rect(s, xs[j], header_y, cw, 0.44, C_VIOLET)
    add_text(s, hdr, xs[j]+0.07, header_y+0.06, cw-0.1, 0.34,
             size=9, bold=True, color=C_WHITE)

for i, (comp, port, db, resp) in enumerate(components):
    row_y = header_y + 0.48 + i * (row_h + 0.03)
    bg_col = BG_MID if i % 2 == 0 else BG_CARD
    total_w = sum(col_widths) + 0.06
    add_rect(s, xs[0], row_y, total_w, row_h, bg_col, BG_MID, 0.5)
    data = [comp, port, db, resp]
    colors_ = [C_VLIGHT, C_TEAL, C_GREEN, C_GRAY]
    for j,(val,cw) in enumerate(zip(data, col_widths)):
        add_text(s, val, xs[j]+0.08, row_y+0.08, cw-0.12, row_h-0.1,
                 size=8.5, color=colors_[j],
                 bold=(j==0), wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDES 5–18 — Sequence Diagram Slides
# ─────────────────────────────────────────────────────────────────────────────
for flow_idx, (title, participants, steps, bullets) in enumerate(FLOWS):
    s = blank(prs)
    set_bg(s, BG_LIGHT)

    flow_num = flow_idx + 1
    slide_header(s, f"Flow {flow_num} — {title}",
                 f"{flow_num} of {len(FLOWS)} flows",
                 dark=False)

    # Render sequence diagram
    png = render_sequence(f"Flow {flow_num} — {title}", participants, steps)
    add_image_bytes(s, png, 0.18, 0.74, 8.35, 6.62)

    # Explanation panel
    bullet_panel(s, f"Flow {flow_num}: {title}", bullets,
                 x=8.72, y=0.72, w=4.38, h=6.62)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 19 — Technology Stack
# ─────────────────────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s, BG_DARK)
slide_header(s, "Technology Stack")

tech_groups = [
    ("Backend", C_VIOLET, [
        ("NestJS v10", "4 independent microservices, TypeScript strict mode"),
        ("TypeORM + SQLite", "better-sqlite3; synchronize:true; per-service DB files"),
        ("Passport JWT + bcrypt", "HS256 tokens; 24h expiry; bcrypt rounds=10"),
        ("Multer", "UUID file uploads; GitHub CDN with /uploads/ fallback"),
    ]),
    ("Frontend", C_TEAL, [
        ("Vue 3 + Vite", "Composition API; <script setup>; Vue Router 4; Pinia"),
        ("Tailwind CSS v4", "CSS-first config in style.css; @theme {} tokens"),
        ("TipTap", "StarterKit + CodeBlockLowlight + Image + Link extensions"),
        ("highlight.js", "Post-render syntax highlighting on .post-content pre code"),
    ]),
    ("AI & ML", C_AMBER, [
        ("GPT-4o + web_search_preview", "Topic discovery; 3-parallel deep research; article writing"),
        ("DALL-E 3/2 → FLUX.1 → Gemini → Unsplash", "5-provider image fallback chain"),
        ("Kokoro AI (local)", "On-device TTS; WAV output; chunked pre-fetch pipeline"),
        ("openai-agents SDK", "Multi-agent rebranding pipeline with handoffs"),
        ("LangGraph 1.2", "Post agent graph; interrupt_before for HITL approval"),
    ]),
    ("Infra & DevOps", C_GREEN, [
        ("Docker + Supervisord", "Single container; 8 processes managed by supervisord"),
        ("Nginx", "Reverse proxy inside Docker; serves frontend assets"),
        ("GitHub Actions", "CI/CD: push → sync secrets → git push to HF Space"),
        ("HuggingFace Spaces", "Docker runtime; mishrabp-myblogs.hf.space"),
        ("node-cron + python-cron", "Sat+Sun 3am UTC automated post generation"),
    ]),
]

col_w = 3.1
for gi, (group_name, accent, items) in enumerate(tech_groups):
    gx = 0.22 + gi * (col_w + 0.1)
    add_rect(s, gx, 0.82, col_w, 0.44, accent)
    add_text(s, group_name, gx+0.1, 0.87, col_w-0.15, 0.35,
             size=10.5, bold=True, color=C_WHITE)
    for ii, (tech, note) in enumerate(items):
        iy = 1.32 + ii * 1.18
        add_rect(s, gx, iy, col_w, 1.12, BG_CARD, accent, 0.8)
        add_text(s, tech, gx+0.12, iy+0.08, col_w-0.2, 0.38,
                 size=9.5, bold=True, color=C_VLIGHT)
        add_text(s, note, gx+0.12, iy+0.5, col_w-0.2, 0.56,
                 size=8, color=C_GRAY, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 20 — Future Improvements
# ─────────────────────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s, BG_DARK)
slide_header(s, "Future Improvements — Roadmap Ideation")

improvements = [
    ("🔍", "Semantic Search",       C_VIOLET,
     "Replace SQL LIKE with dense vector search (sqlite-vec or pgvector). "
     "Embed posts at publish time; query by cosine similarity for meaning-based results."),
    ("🌐", "Multi-Language AI",      C_TEAL,
     "Auto-translate every published post to ES, FR, DE, ZH via AI. "
     "Serve language versions based on Accept-Language; add navbar language switcher."),
    ("📣", "Social Auto-Posting",    C_ROSE,
     "On publish, AI generates platform-optimised captions and auto-posts to LinkedIn, "
     "Twitter/X, and Mastodon. Track social referral traffic back per post."),
    ("🤝", "AI Comment Moderation",  C_AMBER,
     "Replace manual approve/reject with an AI sentiment + spam classifier. "
     "Auto-approve positive comments; auto-reject spam; flag only edge cases for humans."),
    ("📊", "Analytics Microservice", C_GREEN,
     "Beyond view counts: session duration, scroll depth, and read-completion rates. "
     "Lightweight pixel tracker; no third-party dependency; charts in admin dashboard."),
    ("🗓️", "Post Scheduling Queue",  C_VLIGHT,
     "Draft posts with a scheduled publish date. AI agent queues multiple posts and "
     "drip-publishes them. Calendar view in admin panel shows the upcoming schedule."),
    ("🗣️", "Podcast / Audio Feed",   C_AMBER,
     "Convert popular posts to podcast episodes with intro jingles, chapter markers, "
     "and an RSS feed compatible with Apple Podcasts and Spotify."),
    ("🐘", "PostgreSQL Migration",   C_TEAL,
     "Upgrade from SQLite to PostgreSQL for concurrent writes, connection pooling, "
     "pg_trgm fuzzy search, and production-grade reliability. Keep SQLite for dev."),
    ("🖼️", "AI Alt-Text Generation", C_GREEN,
     "On every image upload, call a vision model to auto-generate descriptive alt text. "
     "Stored in media.db; surfaced in the post editor; improves ADA and SEO automatically."),
]

cols = 3
for idx, (icon, title_, accent, desc) in enumerate(improvements):
    col = idx % cols
    row = idx // cols
    cx = 0.22 + col * 4.33
    cy = 0.85 + row * 2.12
    add_rect(s, cx, cy, 4.18, 2.0, BG_CARD, accent, 1)
    add_rect(s, cx, cy, 4.18, 0.46, BG_MID)
    add_text(s, f"{icon}  {title_}", cx+0.12, cy+0.06, 3.9, 0.38,
             size=10.5, bold=True, color=accent)
    add_text(s, desc, cx+0.12, cy+0.56, 3.94, 1.36,
             size=8.5, color=C_GRAY, wrap=True)

add_text(s, "★  These are AI-generated improvement ideations — please tune and prioritise as needed",
         0.22, 7.08, SW-0.4, 0.34,
         size=8.5, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 21 — End
# ─────────────────────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s, BG_DARK)
add_rect(s, 0, 0, 0.22, SH, C_VIOLET)
add_rect(s, 0.22, 0, 0.06, SH, BG_MID)

add_text(s, "Thank You", 1.0, 1.6, 10, 1.3,
         size=52, bold=True, color=C_VLIGHT)
add_rect(s, 1.0, 3.0, 6.5, 0.05, C_VIOLET)
add_text(s, "mishrabp-myblogs.hf.space", 1.0, 3.18, 10, 0.6,
         size=16, color=C_WHITE)
add_text(s, "Bibhu Mishra · Principal Cloud / DevOps Architect · LTIMindtree",
         1.0, 3.88, 10, 0.5,
         size=11, italic=True, color=C_GRAY)
add_text(s, "linkedin.com/in/mishrab",
         1.0, 4.5, 10, 0.45,
         size=10, color=C_TEAL)

# Slide count
add_text(s, f"{len(prs.slides)} slides  ·  14 sequence diagrams  ·  June 2026",
         1.0, 6.9, SW-1.2, 0.42,
         size=9, italic=True, color=C_GRAY, align=PP_ALIGN.LEFT)


# ── Save ───────────────────────────────────────────────────────────────────────
OUTPUT = os.path.join(os.path.dirname(__file__), '..', 'meridian_sysdesign.pptx')
prs.save(OUTPUT)
print(f"✅  PPTX written to {os.path.abspath(OUTPUT)}")
print(f"    Slides: {len(prs.slides)}")
